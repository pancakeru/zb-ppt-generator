from pptx import Presentation
from pptx.util import Inches, Pt
import os
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

template = Presentation("/Users/zoe.chow/Desktop/ai-test/Riot Games PPT Template.pptx")

from datetime import datetime, timedelta

def get_week_range():
    today = datetime.today()
    start_of_week = today - timedelta(days=today.weekday() + 1 if today.weekday() != 6 else 0)
    end_of_week = start_of_week + timedelta(days=6)
    return start_of_week.strftime('%m%d'), end_of_week.strftime('%m%d')

start_str, end_str = get_week_range()
date_range_str = f"{start_str}-{end_str}" 
domestic_labels = ["B站","抖音", "微信群", "公众号", "小红书"]
international_labels = ["X/Twitter", "Discord", "Facebook", "YouTube", "Instagram", "Reddit"]

#for i, layout in enumerate(template.slide_layouts):
    #print(f"{i}: {layout.name}")

def make_ppt(data_list, yt_data, yt_keywords, bb_data, bb_keywords, output_path=f"/Users/zoe.chow/Desktop/ai-test/25{date_range_str}符文战场周报.pptx"):
    prs = Presentation()
    prs.slide_width = template.slide_width
    prs.slide_height = template.slide_height
    title_slide_layout = template.slide_layouts[1]  # Title slide
    template_slide = template.slides[4]

    # === Title Slide ===
    slide = prs.slides.add_slide(title_slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_text= f"{date_range_str}符文战场周报\n"
    subtitle_text = "自动生成"

    title_shape.text = title_text
    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_run = title_paragraph.runs[0]

    title_paragraph.alignment = PP_ALIGN.CENTER  
    title_run.font.size = Pt(44)
    title_run.font.name = 'Arial Black'
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(255, 255, 255) 

    # --- Style the subtitle ---
    subtitle_shape.text = subtitle_text
    subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
    subtitle_run = subtitle_paragraph.runs[0]

    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_run.font.size = Pt(23)
    subtitle_run.font.name = 'Georgia'
    subtitle_run.font.color.rgb = RGBColor(255, 0, 0)

    print("Title slide added")

    # ---Section header
    add_section_slide(prs, "01 社群讨论", ["国内玩家", "国外玩家", "总结"])
    print("Section 1 added")

    # --- Community slides
    add_social_summary_slide(
        prs,
        "国内玩家：",
        domestic_labels
    )
    print("Domestic slide added")
    add_social_summary_slide(
        prs,
        "国外玩家：",
        international_labels
    )
    print("International slide added")
    add_community_summary_slide(
        prs,
        left_lines=("IP粉丝...", "卡牌玩家..."),
        cn_title="国内玩家",
        cn_sub="玩家...",
        intl_title="国外玩家",
        intl_sub="玩家...",
    )
    print("Community Summary slide added")

    # ---- YouTube and BiliBili
    add_section_slide(prs, "02 视频热度", ["B站", "YouTube"])
    print("Section 2 added")
    add_data_slide(bb_data, bb_keywords, prs, "符文战场近10天热门B站视频排行榜")
    print("BiliBili slide added")
    add_data_slide(yt_data, yt_keywords, prs, "符文战场近10天热门YouTube视频排行榜")
    print("YouTube slide added")

    add_section_slide(prs, "03 竞品动向", ["宝可梦", "航海王", "数码宝贝", "高达卡牌"])
    print("Section 3 added")

    # === Content Slides ===
    for item in data_list:
        fill_from_custom_slide(prs, template_slide, item)

    # ---- Products Summary
    create_top3_slide(prs, ["宝可梦...", "航海王...", "数码宝贝..."])
    print("Product summary slide added")

    # === Save PPT ===
    prs.save(output_path)
    print(f"✅ PPT saved to {output_path}")

import requests
from io import BytesIO

def download_and_insert_image(image_url, slide, left, top, width=None, height=None):
    try:
        response = requests.get(image_url)
        if response.status_code == 200:
            img_stream = BytesIO(response.content)
            slide.shapes.add_picture(img_stream, left, top, width, height)
        else:
            print(f"⚠️ Failed to download image from {image_url}")
    except Exception as e:
        print(f"❌ Error downloading image: {e}")

import requests
from io import BytesIO

def fill_from_custom_slide(prs, template_slide, data):
    # Duplicate the layout
    slide_layout = template_slide.slide_layout
    slide = prs.slides.add_slide(slide_layout)

    # Copy the background and shapes
    for shape in template_slide.shapes:
        if shape.shape_type == 13:  # Picture
            img_stream = BytesIO(shape.image.blob)
            slide.shapes.add_picture(img_stream, shape.left, shape.top, shape.width, shape.height)
        elif shape.has_text_frame:
            new_shape = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            new_shape.text_frame.word_wrap = True 

            # Decide what to fill based on existing placeholder content
            original_text = shape.text.strip().lower()
            if "insert image" in original_text and data.get("image"):
                try:
                    response = requests.get(data["image"])
                    if response.status_code == 200:
                        try:
                            img_stream = BytesIO(response.content)
                            slide.shapes.add_picture(img_stream, shape.left, shape.top, shape.width, shape.height)
                        except ValueError as e:
                            print(f"⚠️ Standard insert failed: {e}, trying webp conversion...")
                            converted_path = convert_webp_to_png(response.content)
                            if converted_path:
                                slide.shapes.add_picture(converted_path, shape.left, shape.top, shape.width, shape.height)
                                os.remove(converted_path)  # Cleanup temp png
                except Exception as e:
                    print(f"❌ Unexpected error: {e}")
            elif "subtitle" in original_text:
                new_text_lines = [
                    " ",
                    f"发布日期：{data['date']}",
                    f"类别：{data['type']}",
                    f"{data['info']}\n",
                   ("链接：", data['link'])
                ]
                new_shape.text_frame.clear()
                new_shape.text_frame.word_wrap = True

                for i, line in enumerate(new_text_lines):
                    if i == 0:
                        para = new_shape.text_frame.paragraphs[0]
                    else:
                        para = new_shape.text_frame.add_paragraph()

                    if isinstance(line, tuple):
                        label, url = line
                        run_label = para.add_run()
                        run_label.text = label
                        run_label.font.color.rgb = RGBColor(255, 0, 0)
                        run_label.font.size = Pt(11)

                        run_link = para.add_run()
                        run_link.text = url
                        run_link.hyperlink.address = url
                        run_link.font.color.rgb = RGBColor(255, 0, 0)
                        run_link.font.size = Pt(11)
                    else:
                        para.text = line
                        para.font.size = Pt(14 if i == 1 else 11)
                        para.font.color.rgb = RGBColor(255, 255, 255)
            elif "sample image layout" in original_text:
                para = new_shape.text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER 
                para.text = f"{data['name']}{data['flavor']}"

                run = para.runs[0]
                run.font.size = Pt(28) 
                run.font.bold = True
            else:
                new_shape.text_frame.text = ""
    print("slide added")
    return slide

def add_section_slide(prs, header_text, side_list):
    layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(layout)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        sp = shape
        sp.text_frame.clear()

    # Left section title
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5), Inches(1))
    left_tf = left_box.text_frame
    left_p = left_tf.paragraphs[0]
    left_p.text = header_text
    left_p.font.size = Pt(44)
    left_p.font.color.rgb = RGBColor(255, 255, 255)
    left_p.font.bold = True

    # Right vertical list
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(2.5), Inches(3))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True

    for word in side_list:
        p = right_tf.add_paragraph()
        p.text = f"{word}\n\n"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

    return slide

from PIL import Image
import requests
from io import BytesIO
import os

def convert_webp_to_png(webp_bytes, output_path="temp_image.png"):
    try:
        img = Image.open(BytesIO(webp_bytes))
        if img.format == "WEBP":
            img.save(output_path, "PNG")
            return output_path
    except Exception as e:
        print(f"❌ Conversion failed: {e}")
    return None

def add_data_slide(videos, keywords, prs, header_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add left box (Top Videos)
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7), Inches(5))
    tf_left = left_box.text_frame
    tf_left.clear()

    title = tf_left.paragraphs[0]
    title.text = header_text
    title.runs[0].font.size = Pt(26)
    title.runs[0].font.bold = True

    for video in videos:
    # Title - white
        title_para = tf_left.add_paragraph()
        title_para.text = f"\n• {video['title']}\n{video['channel']} | {video['views']} 次观看"
        title_para.font.size = Pt(13)
        title_para.font.bold = False
        title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Link - red
        link_para = tf_left.add_paragraph()
        link_run = link_para.add_run()
        link_run.text = video['url']
        link_para.font.size = Pt(12)
        link_para.font.color.rgb = RGBColor(255, 0, 0)
        link_run.hyperlink.address = video['url']

    # Add right box (Top Keywords)
    right_box = slide.shapes.add_textbox(Inches(7.5), Inches(1), Inches(3), Inches(5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    tf_right.clear()

    ktitle = tf_right.paragraphs[0]
    ktitle.text = "关键词:"
    ktitle.runs[0].font.size = Pt(16)
    ktitle.runs[0].font.bold = True

    for keyword, count in keywords:
        para = tf_right.add_paragraph()
        para.text = f"- {keyword} - {count}"
        para.level = 1
        para.runs[0].font.size = Pt(13)

    return slide

def create_top3_slide(prs, items, title_text="动向速览"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = 1  # Center

    # Position settings
    left_num = Inches(1)
    left_text = Inches(1.6)
    top = Inches(1.5)
    num_width = Inches(0.5)
    text_width = Inches(6.5)
    height = Inches(1)
    spacing = Inches(1.4)

    for i, item in enumerate(items):
        y = top + spacing * i

        # Number box (white background)
        num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_num, y, num_width, height)
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        num_box.line.fill.background()
        num_frame = num_box.text_frame
        num_frame.text = str(i + 1)
        num_frame.paragraphs[0].font.size = Pt(36)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0) 
        num_frame.paragraphs[0].alignment = 1 

        # Title box (gray background)
        text_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_text, y, text_width, height)
        text_box.fill.solid()
        text_box.fill.fore_color.rgb = RGBColor(60, 60, 60)  # dark gray
        text_box.line.fill.background()
        text_frame = text_box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = item
        run.font.size = Pt(21)
        run.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

def add_social_summary_slide(
    prs: Presentation,
    title: str,
    labels: list[str],
    *,
    title_left=Inches(0.6), title_top=Inches(0.6),
    title_w=Inches(9), title_h=Inches(1.0),
    rows_left=Inches(0.6), rows_top=Inches(1.6),
    rows_w=Inches(9), rows_h=Inches(5.5),
    title_size=36, label_size=24, value_size=18,
    value_placeholder: str = ""   # set to "—" if you want a visible dash
):
    """Slide with a big title and a single textbox of rows like '标签：  ' (blank values)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # --- Title ---
    tb_title = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf_t = tb_title.text_frame
    tf_t.clear()
    p = tf_t.paragraphs[0]
    p.text = title or ""
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0] if p.runs else p.add_run()
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    # --- Rows (one paragraph per label) ---
    tb_rows = slide.shapes.add_textbox(rows_left, rows_top, rows_w, rows_h)
    tf_r = tb_rows.text_frame
    tf_r.clear()

    for i, label in enumerate(labels):
        para = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        para.alignment = PP_ALIGN.LEFT

        run_label = para.add_run()
        run_label.text = f"{label}："
        run_label.font.size = Pt(label_size)
        run_label.font.bold = True
        run_label.font.color.rgb = RGBColor(255, 255, 255)

        run_space = para.add_run(); run_space.text = " "

        run_val = para.add_run()
        run_val.text = value_placeholder  # stays empty unless you set a placeholder
        run_val.font.size = Pt(value_size)
        run_val.font.bold = False
        run_val.font.color.rgb = RGBColor(255, 255, 255)

    return slide

def add_community_summary_slide(
    prs,
    *,
    left_lines=("IP粉丝…", "卡牌玩家…"),
    cn_title="国内玩家",
    cn_sub="玩家…",
    intl_title="国外玩家",
    intl_sub="玩家…",
    panel_fill=RGBColor(0, 0, 0),
    ):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ---- Left slanted panel (trapezoid) ----
    # size/position: cover ~40% width, full height; slant via adjustment
    left = Inches(0.0)
    top = Inches(0.0)
    width = Inches(5.2)
    height = Inches(7.5)  # typical 16:9 slide ~7.5" tall

    shp = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, left, top, width, height)
    # Make the right edge slanted by pulling top width inward (0.0..1.0)
    try:
        shp.adjustments[0] = 0.35         # bigger = stronger slant
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = panel_fill
    shp.fill.fore_color.brightness = 0  # solid black
    shp.line.fill.background()          # no border

    # ---- Left panel text (two lines) ----
    # First line
    tb1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(3.5), Inches(0.9))
    tf1 = tb1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = left_lines[0] if len(left_lines) > 0 else ""
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.runs[0]
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(255, 255, 255)

    # Second line
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(3.8), Inches(0.9))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = left_lines[1] if len(left_lines) > 1 else ""
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.runs[0]
    r2.font.size = Pt(28)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(255, 255, 255)

    # ---- Right column blocks ----
    # Common styling helpers
    def add_block(title, sub, top_in):
        # Title (bigger)
        tb_t = slide.shapes.add_textbox(Inches(6.2), top_in, Inches(3.8), Inches(0.9))
        tf_t = tb_t.text_frame; tf_t.clear()
        pt = tf_t.paragraphs[0]
        pt.text = title
        pt.alignment = PP_ALIGN.LEFT
        rt = pt.runs[0]
        rt.font.size = Pt(32)
        rt.font.bold = True
        rt.font.color.rgb = RGBColor(255, 255, 255)

        # Sub (smaller, lighter)
        tb_s = slide.shapes.add_textbox(Inches(6.2), top_in + Inches(0.8), Inches(3.8), Inches(0.7))
        tf_s = tb_s.text_frame; tf_s.clear()
        ps = tf_s.paragraphs[0]
        ps.text = sub
        ps.alignment = PP_ALIGN.LEFT
        rs = ps.runs[0]
        rs.font.size = Pt(20)
        rs.font.bold = False
        rs.font.color.rgb = RGBColor(200, 200, 200)

    add_block(cn_title,   cn_sub,   Inches(1.6))
    add_block(intl_title, intl_sub, Inches(3.8))

    return slide